from pptx import Presentation

def extract_ppt(path):
    prs = Presentation(path)
    return " ".join(shape.text for s in prs.slides for shape in s.shapes if hasattr(shape,"text"))
